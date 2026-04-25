#!/usr/bin/env python3
"""Build demo.pptx from a run directory containing steps_described.json + screenshots/.

Usage:
    python3 output/pptGenerator.py <run_dir>

The script reads <run_dir>/steps_described.json (preferred) or steps.json,
and emits <run_dir>/demo.pptx with one cover slide and one slide per step.

Each step slide contains:
    - Title (LLM-generated, falls back to step label)
    - Description (LLM-generated, falls back to step label)
    - The step's screenshot, scaled to fit the bottom area while preserving
      aspect ratio.

Dependencies:
    pip install python-pptx Pillow
"""
from __future__ import annotations

import json
import sys
from datetime import datetime
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt
except ModuleNotFoundError:
    sys.stderr.write(
        "error: python-pptx not installed.\n"
        "       run:  pip install python-pptx Pillow\n"
    )
    sys.exit(2)

try:
    from PIL import Image
except ModuleNotFoundError:
    Image = None  # We can still emit slides without exact aspect ratios.


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.5)
TITLE_H = Inches(0.7)
DESC_H = Inches(0.9)

ACCENT = RGBColor(0x0F, 0x76, 0xFF)
TEXT = RGBColor(0x1F, 0x2A, 0x44)
MUTED = RGBColor(0x55, 0x66, 0x80)
PANEL_BG = RGBColor(0xF5, 0xF7, 0xFB)


def build(run_dir: Path) -> Path:
    manifest = load_manifest(run_dir)
    flow_name = manifest.get("flow_name") or run_dir.name
    generated_at = manifest.get("generated_at") or datetime.utcnow().isoformat() + "Z"
    steps = manifest.get("steps") or []

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    add_cover_slide(prs, blank, flow_name, generated_at, len(steps))

    for step in steps:
        add_step_slide(prs, blank, step, run_dir)

    out = run_dir / "demo.pptx"
    prs.save(str(out))
    return out


def load_manifest(run_dir: Path) -> dict:
    described = run_dir / "steps_described.json"
    if described.exists():
        return json.loads(described.read_text())
    plain = run_dir / "steps.json"
    if plain.exists():
        return json.loads(plain.read_text())
    raise SystemExit(f"no manifest found in {run_dir} (expected steps_described.json or steps.json)")


def add_cover_slide(prs, layout, flow_name: str, generated_at: str, step_count: int) -> None:
    slide = prs.slides.add_slide(layout)
    add_band(slide, top=0, height=Inches(2.2), color=ACCENT)
    add_text_box(
        slide,
        left=MARGIN,
        top=Inches(0.6),
        width=SLIDE_W - 2 * MARGIN,
        height=Inches(0.7),
        text="Quest Data Modeler",
        font_size=Pt(20),
        color=RGBColor(0xFF, 0xFF, 0xFF),
        bold=False,
    )
    add_text_box(
        slide,
        left=MARGIN,
        top=Inches(1.15),
        width=SLIDE_W - 2 * MARGIN,
        height=Inches(0.9),
        text="Automated Demo Walkthrough",
        font_size=Pt(36),
        color=RGBColor(0xFF, 0xFF, 0xFF),
        bold=True,
    )
    pretty_when = format_timestamp(generated_at)
    add_text_box(
        slide,
        left=MARGIN,
        top=Inches(2.6),
        width=SLIDE_W - 2 * MARGIN,
        height=Inches(0.6),
        text=flow_name,
        font_size=Pt(22),
        color=TEXT,
        bold=True,
    )
    add_text_box(
        slide,
        left=MARGIN,
        top=Inches(3.2),
        width=SLIDE_W - 2 * MARGIN,
        height=Inches(0.5),
        text=f"{step_count} steps · generated {pretty_when}",
        font_size=Pt(14),
        color=MUTED,
        bold=False,
    )


def add_step_slide(prs, layout, step: dict, run_dir: Path) -> None:
    index = step.get("index", 0)
    title = step.get("title") or step.get("label") or f"Step {index}"
    description = step.get("description") or step.get("label") or ""
    screenshot = step.get("screenshot")

    slide = prs.slides.add_slide(layout)

    add_band(slide, top=0, height=Inches(0.18), color=ACCENT)

    add_text_box(
        slide,
        left=MARGIN,
        top=Inches(0.32),
        width=Inches(0.9),
        height=Inches(0.5),
        text=f"Step {index:02d}",
        font_size=Pt(12),
        color=ACCENT,
        bold=True,
    )
    add_text_box(
        slide,
        left=Inches(1.4),
        top=Inches(0.28),
        width=SLIDE_W - Inches(1.9),
        height=TITLE_H,
        text=title,
        font_size=Pt(24),
        color=TEXT,
        bold=True,
    )
    add_text_box(
        slide,
        left=MARGIN,
        top=Inches(1.0),
        width=SLIDE_W - 2 * MARGIN,
        height=DESC_H,
        text=description,
        font_size=Pt(14),
        color=MUTED,
        bold=False,
    )

    image_top = Inches(2.0)
    image_bottom = SLIDE_H - Inches(0.4)
    image_max_w = SLIDE_W - 2 * MARGIN
    image_max_h = image_bottom - image_top

    if screenshot:
        path = Path(screenshot)
        if not path.is_absolute():
            path = run_dir / path
        if path.exists():
            place_image(slide, str(path), image_top, image_max_w, image_max_h)
        else:
            add_text_box(
                slide,
                left=MARGIN,
                top=image_top,
                width=image_max_w,
                height=Inches(0.5),
                text=f"(screenshot not found: {path})",
                font_size=Pt(12),
                color=MUTED,
                bold=False,
            )


def place_image(slide, image_path: str, top, max_w, max_h) -> None:
    width, height = max_w, max_h
    if Image is not None:
        try:
            with Image.open(image_path) as im:
                iw, ih = im.size
            scale = min(max_w / iw, max_h / ih)
            width = int(iw * scale)
            height = int(ih * scale)
        except Exception:
            pass
    left = int((SLIDE_W - width) / 2)
    slide.shapes.add_picture(image_path, left, top, width=width, height=height)


def add_band(slide, top, height, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, SLIDE_W, height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text: str, *, font_size, color: RGBColor, bold: bool) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color


def format_timestamp(s: str) -> str:
    try:
        dt = datetime.fromisoformat(s.replace("Z", "+00:00"))
        return dt.strftime("%b %d, %Y · %H:%M UTC")
    except Exception:
        return s


def main() -> None:
    if len(sys.argv) < 2:
        sys.stderr.write("usage: python3 output/pptGenerator.py <run_dir>\n")
        sys.exit(1)
    run_dir = Path(sys.argv[1]).resolve()
    if not run_dir.is_dir():
        sys.stderr.write(f"error: not a directory: {run_dir}\n")
        sys.exit(1)
    out = build(run_dir)
    print(str(out))


if __name__ == "__main__":
    main()
