"""
Generate an executive-friendly PPTX UX audit of Quest Data Modeler.

Lenses: Hick's Law (reduce choices), Fitts's Law (make actions easy),
Jakob's Law (keep things familiar). Findings reference real screenshots
from the demo run.

Output: output/qdm_ux_audit.pptx
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
RUN = ROOT / "output" / "quest_data_modeler_end_to_end_walkthrough_2026-04-26T11-46-09-041Z"
SCREENS = RUN / "screenshots"
LOGO = ROOT / "assets" / "quest_logo.png"
OUT = ROOT / "output" / "qdm_ux_audit.pptx"

# Quest brand palette
NAVY = RGBColor(0x0E, 0x1F, 0x4D)
GOLD = RGBColor(0xF5, 0xB7, 0x00)
INK = RGBColor(0x1B, 0x1F, 0x2A)
MUTED = RGBColor(0x6B, 0x72, 0x82)
ACCENT = RGBColor(0x14, 0x6C, 0xF1)
CARD = RGBColor(0xF4, 0xF6, 0xFB)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Slide canvas: 13.333 x 7.5 inches (16:9 widescreen)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def add_textbox(slide, left, top, width, height, text, *, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return box


def add_bullets(slide, left, top, width, height, items, *, size=14, color=INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # bullet symbol
        bullet = p.add_run()
        bullet.text = "•  "
        bullet.font.size = Pt(size)
        bullet.font.color.rgb = ACCENT
        bullet.font.bold = True
        bullet.font.name = "Calibri"
        # text
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
        p.space_after = Pt(8)
    return box


def add_rect(slide, left, top, width, height, fill, line=None, line_color=None):
    from pptx.enum.shapes import MSO_SHAPE
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line_color or NAVY
        rect.line.width = line
    rect.shadow.inherit = False
    return rect


def add_logo(slide, left, top, height):
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), left, top, height=height)


def header_band(slide, kicker: str, title: str):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.55), NAVY)
    add_logo(slide, Inches(0.4), Inches(0.08), Inches(0.4))
    add_textbox(slide, Inches(1.1), Inches(0.08), Inches(8), Inches(0.4), kicker, size=12, bold=True, color=GOLD)
    add_textbox(slide, Inches(8.5), Inches(0.13), Inches(4.6), Inches(0.4), title, size=11, color=WHITE, align=PP_ALIGN.RIGHT)


# ---------- Slides ----------

def slide_cover(prs):
    s = add_blank(prs)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, NAVY)
    # Diagonal accent
    add_rect(s, Inches(0), Inches(6.8), SLIDE_W, Inches(0.7), GOLD)
    add_logo(s, Inches(0.6), Inches(0.6), Inches(0.85))
    add_textbox(s, Inches(0.6), Inches(2.2), Inches(11), Inches(0.6),
                "Quest Data Modeler", size=40, bold=True, color=WHITE)
    add_textbox(s, Inches(0.6), Inches(3.0), Inches(11), Inches(0.9),
                "UX Audit", size=64, bold=True, color=WHITE)
    add_textbox(s, Inches(0.6), Inches(4.4), Inches(11), Inches(0.5),
                "Hick's Law  ·  Fitts's Law  ·  Jakob's Law", size=22, color=GOLD)
    add_textbox(s, Inches(0.6), Inches(5.2), Inches(11), Inches(0.5),
                "Findings from the end-to-end demo walkthrough", size=18, color=WHITE)
    add_textbox(s, Inches(0.6), Inches(6.95), Inches(11), Inches(0.4),
                "Prepared for stakeholder review · April 2026", size=12, color=NAVY, bold=True)


def slide_exec_summary(prs):
    s = add_blank(prs)
    header_band(s, "EXECUTIVE SUMMARY", "Quest Data Modeler — UX Audit")
    add_textbox(s, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.6),
                "Three usability laws · 11 findings · 6 quick wins", size=24, bold=True, color=NAVY)

    # Three columns — one per law
    cards = [
        (NAVY, "Hick's Law", "Reduce choices",
         "16 unlabeled icons in the canvas toolbar slow first-time users.",
         "Group + label."),
        (ACCENT, "Fitts's Law", "Make actions easy",
         "28-px icons and corner-pinned controls force pixel hunts.",
         "Bigger targets, predictable placement."),
        (GOLD, "Jakob's Law", "Keep things familiar",
         "AI panel mirrors ChatGPT — but disabled controls and silent failures break the mental model.",
         "Hide disabled UI; commit user choices visibly."),
    ]
    col_w = Inches(4.0)
    gap = Inches(0.25)
    start_x = Inches(0.5)
    top = Inches(1.7)
    for i, (color, law, sub, finding, fix) in enumerate(cards):
        x = Emu(int(start_x) + i * (int(col_w) + int(gap)))
        add_rect(s, x, top, col_w, Inches(0.6), color)
        add_textbox(s, Emu(int(x) + Inches(0.2)), Emu(int(top) + Inches(0.1)), col_w, Inches(0.5),
                    law, size=20, bold=True, color=WHITE)
        # body card
        body_top = Emu(int(top) + int(Inches(0.6)))
        add_rect(s, x, body_top, col_w, Inches(3.6), CARD)
        add_textbox(s, Emu(int(x) + Inches(0.2)), Emu(int(body_top) + Inches(0.15)), col_w, Inches(0.4),
                    sub, size=14, bold=True, color=color)
        add_textbox(s, Emu(int(x) + Inches(0.2)), Emu(int(body_top) + Inches(0.65)), Inches(3.6), Inches(1.6),
                    "What we saw:", size=11, bold=True, color=MUTED)
        add_textbox(s, Emu(int(x) + Inches(0.2)), Emu(int(body_top) + Inches(0.95)), Inches(3.6), Inches(1.8),
                    finding, size=13, color=INK)
        add_textbox(s, Emu(int(x) + Inches(0.2)), Emu(int(body_top) + Inches(2.4)), Inches(3.6), Inches(0.4),
                    "Recommendation:", size=11, bold=True, color=MUTED)
        add_textbox(s, Emu(int(x) + Inches(0.2)), Emu(int(body_top) + Inches(2.7)), Inches(3.6), Inches(0.8),
                    fix, size=13, bold=True, color=GREEN)

    add_textbox(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
                "Top opportunity → Replace the unlabeled bottom toolbar with grouped, labeled controls. Highest impact, lowest effort.",
                size=14, bold=True, color=NAVY)


def slide_methodology(prs):
    s = add_blank(prs)
    header_band(s, "METHODOLOGY", "How these findings were captured")
    add_textbox(s, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.6),
                "End-to-end task observation, then evaluated through three classic UX heuristics.",
                size=18, color=INK)

    # Tasks observed
    add_textbox(s, Inches(0.5), Inches(1.7), Inches(6), Inches(0.4),
                "TASKS OBSERVED", size=12, bold=True, color=NAVY)
    add_bullets(s, Inches(0.5), Inches(2.05), Inches(6), Inches(4),
                [
                    "Sign in to Quest Data Modeler",
                    "Generate a model from a natural-language prompt (AI Generator)",
                    "Edit metadata in Quick Editor (rename Customer → Customers)",
                    "Pivot across Entities / Domains / Relationships",
                    "Add a new entity (OrderLine) on the canvas",
                    "Forward-engineer SQL on a Logical/Physical model",
                ], size=14)

    # Lenses
    add_textbox(s, Inches(7.0), Inches(1.7), Inches(6), Inches(0.4),
                "EVALUATION LENSES", size=12, bold=True, color=NAVY)
    lenses = [
        ("Hick's Law", "Time to decide grows with the number of choices presented."),
        ("Fitts's Law", "Time to acquire a target depends on its size and distance."),
        ("Jakob's Law", "Users prefer products that work like the others they already know."),
    ]
    y = 2.05
    for name, defn in lenses:
        add_textbox(s, Inches(7.0), Inches(y), Inches(6), Inches(0.4),
                    name, size=15, bold=True, color=ACCENT)
        add_textbox(s, Inches(7.0), Inches(y + 0.4), Inches(6), Inches(0.6),
                    defn, size=12, color=INK)
        y += 1.1

    add_textbox(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
                "Source: 63-step recorded walkthrough · screenshots referenced throughout this deck.",
                size=11, color=MUTED)


def slide_law_intro(prs, law: str, sub: str, defn: str, screenshot: str | None):
    s = add_blank(prs)
    header_band(s, law.upper(), sub)
    add_textbox(s, Inches(0.5), Inches(0.85), Inches(8), Inches(0.7),
                law, size=44, bold=True, color=NAVY)
    add_textbox(s, Inches(0.5), Inches(1.6), Inches(8), Inches(0.5),
                sub, size=22, color=GOLD)
    add_textbox(s, Inches(0.5), Inches(2.3), Inches(8), Inches(2.5),
                defn, size=16, color=INK)
    if screenshot and (SCREENS / screenshot).exists():
        s.shapes.add_picture(str(SCREENS / screenshot), Inches(8.5), Inches(1.7), width=Inches(4.5))
        add_textbox(s, Inches(8.5), Inches(5.2), Inches(4.5), Inches(0.4),
                    f"Reference: {screenshot}", size=10, color=MUTED, align=PP_ALIGN.CENTER)


def slide_finding(prs, law_color, kicker, title, observation, impact, recommendation, screenshot, severity):
    s = add_blank(prs)
    header_band(s, kicker, title)

    # Severity chip
    sev_color = {"High": RED, "Medium": GOLD, "Low": MUTED}[severity]
    add_rect(s, Inches(0.5), Inches(0.85), Inches(1.4), Inches(0.4), sev_color)
    add_textbox(s, Inches(0.5), Inches(0.9), Inches(1.4), Inches(0.3),
                f"  {severity.upper()}", size=11, bold=True, color=WHITE)

    # Title
    add_textbox(s, Inches(2.05), Inches(0.85), Inches(11), Inches(0.6),
                title, size=24, bold=True, color=NAVY)

    # Left column: text
    add_textbox(s, Inches(0.5), Inches(1.6), Inches(7.0), Inches(0.4),
                "WHAT WE SAW", size=11, bold=True, color=MUTED)
    add_textbox(s, Inches(0.5), Inches(1.95), Inches(7.0), Inches(2.2),
                observation, size=14, color=INK)

    add_textbox(s, Inches(0.5), Inches(3.6), Inches(7.0), Inches(0.4),
                "WHY IT MATTERS", size=11, bold=True, color=MUTED)
    add_textbox(s, Inches(0.5), Inches(3.95), Inches(7.0), Inches(1.6),
                impact, size=14, color=INK)

    add_textbox(s, Inches(0.5), Inches(5.4), Inches(7.0), Inches(0.4),
                "RECOMMENDATION", size=11, bold=True, color=MUTED)
    add_rect(s, Inches(0.5), Inches(5.7), Inches(7.0), Inches(1.5), CARD)
    add_textbox(s, Inches(0.7), Inches(5.85), Inches(6.8), Inches(1.3),
                recommendation, size=14, bold=True, color=GREEN)

    # Right column: screenshot
    if screenshot and (SCREENS / screenshot).exists():
        # Frame
        add_rect(s, Inches(7.85), Inches(1.55), Inches(5.3), Inches(5.0), CARD, line=Emu(2000), line_color=law_color)
        s.shapes.add_picture(str(SCREENS / screenshot), Inches(7.95), Inches(1.65), width=Inches(5.1))
        add_textbox(s, Inches(7.85), Inches(6.7), Inches(5.3), Inches(0.4),
                    f"From: {screenshot}", size=10, color=MUTED, align=PP_ALIGN.CENTER)


def slide_priority_matrix(prs):
    s = add_blank(prs)
    header_band(s, "PRIORITY MATRIX", "What to ship first")
    add_textbox(s, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.5),
                "Impact vs. effort — start in the upper-left.", size=18, color=INK)

    # Matrix area
    mx, my = Inches(1.0), Inches(1.7)
    mw, mh = Inches(11.3), Inches(5.0)
    add_rect(s, mx, my, mw, mh, CARD)
    # Quadrant lines
    add_rect(s, Emu(int(mx) + int(mw) // 2), my, Emu(2000), mh, MUTED)
    add_rect(s, mx, Emu(int(my) + int(mh) // 2), mw, Emu(2000), MUTED)

    # Axis labels
    add_textbox(s, mx, Emu(int(my) + int(mh) + 30000), mw, Inches(0.4),
                "Effort →", size=12, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, Emu(int(mx) - int(Inches(0.5))), my, Inches(0.4), mh,
                "Impact →", size=12, bold=True, color=MUTED)
    add_textbox(s, Emu(int(mx) + Inches(0.2)), Emu(int(my) + Inches(0.1)), Inches(2), Inches(0.3),
                "Quick wins", size=12, bold=True, color=GREEN)
    add_textbox(s, Emu(int(mx) + int(mw) // 2 + Inches(0.2)), Emu(int(my) + Inches(0.1)), Inches(3), Inches(0.3),
                "Strategic bets", size=12, bold=True, color=NAVY)
    add_textbox(s, Emu(int(mx) + Inches(0.2)), Emu(int(my) + int(mh) // 2 + Inches(0.1)), Inches(3), Inches(0.3),
                "Nice to have", size=12, bold=True, color=MUTED)
    add_textbox(s, Emu(int(mx) + int(mw) // 2 + Inches(0.2)), Emu(int(my) + int(mh) // 2 + Inches(0.1)), Inches(3), Inches(0.3),
                "Defer", size=12, bold=True, color=MUTED)

    # Items: (label, x_frac, y_frac, color)
    items = [
        # Quick wins (low effort, high impact) — upper-left
        ("Tooltips on bottom toolbar", 0.10, 0.20, GREEN),
        ("Hide disabled type dropdown", 0.18, 0.32, GREEN),
        ("Single-click to open a model", 0.28, 0.18, GREEN),
        ("Label AI panel icons", 0.16, 0.45, GREEN),
        # Strategic bets (high effort, high impact) — upper-right
        ("Fix LP-chip → AI-output handoff (bug)", 0.62, 0.18, NAVY),
        ("Re-group canvas toolbar", 0.78, 0.30, NAVY),
        ("Inline AI controls (no popover)", 0.70, 0.40, NAVY),
        # Low impact / low effort — lower-left
        ("Resize wizard close affordance", 0.20, 0.70, MUTED),
        # Defer (high effort, low impact) — lower-right
        ("Move toolbar to top", 0.80, 0.78, MUTED),
    ]
    for label, fx, fy, color in items:
        cx = Emu(int(mx) + int(int(mw) * fx))
        cy = Emu(int(my) + int(int(mh) * fy))
        from pptx.enum.shapes import MSO_SHAPE
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, Inches(0.22), Inches(0.22))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.color.rgb = WHITE
        dot.line.width = Emu(8000)
        add_textbox(s, Emu(int(cx) + int(Inches(0.32))), Emu(int(cy) - int(Inches(0.04))),
                    Inches(4.0), Inches(0.4), label, size=11, color=INK, bold=True)


def slide_roadmap(prs):
    s = add_blank(prs)
    header_band(s, "ROADMAP", "Sequencing the work")
    add_textbox(s, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.5),
                "Three phases. Each phase lifts a measurable usability metric.", size=18, color=INK)

    phases = [
        ("Phase 1 — Quick Wins", "1 sprint", GREEN, [
            "Tooltips on every icon in the canvas toolbar",
            "Replace dblclick-to-open with single click on model name",
            "Hide the disabled type dropdown when not editable",
            "Add labels next to AI panel paperclip + sliders icons",
        ]),
        ("Phase 2 — Trust the User's Choices", "2 sprints", ACCENT, [
            "Persist Logical/Physical + Database chips through AI generation",
            "Move AI Type/DB controls out of the popover, inline below the prompt",
            "Add a confirmation banner showing the type the AI will produce",
            "Surface 'Forward Engineering' state when model can't support it",
        ]),
        ("Phase 3 — Reframe the Canvas", "1 quarter", NAVY, [
            "Group the canvas bottom toolbar into 3 labeled clusters: Build · View · Tools",
            "Reduce overview Quick Actions to one hero CTA + secondary list",
            "Re-evaluate top vs bottom toolbar placement with usability tests",
        ]),
    ]
    y = 1.7
    for title, dur, color, items in phases:
        add_rect(s, Inches(0.5), Inches(y), Inches(12.3), Inches(1.55), CARD)
        add_rect(s, Inches(0.5), Inches(y), Inches(0.18), Inches(1.55), color)
        add_textbox(s, Inches(0.85), Inches(y + 0.1), Inches(8), Inches(0.4),
                    title, size=18, bold=True, color=NAVY)
        add_textbox(s, Inches(0.85), Inches(y + 0.5), Inches(8), Inches(0.3),
                    dur, size=12, bold=True, color=color)
        # bullets in two columns
        col1 = items[: (len(items) + 1) // 2]
        col2 = items[(len(items) + 1) // 2:]
        add_bullets(s, Inches(5.5), Inches(y + 0.15), Inches(3.7), Inches(1.4), col1, size=11)
        add_bullets(s, Inches(9.3), Inches(y + 0.15), Inches(3.7), Inches(1.4), col2, size=11)
        y += 1.75


def slide_appendix(prs):
    s = add_blank(prs)
    header_band(s, "APPENDIX", "Sources & next steps")
    add_textbox(s, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.5),
                "Source material & how to validate", size=22, bold=True, color=NAVY)

    add_textbox(s, Inches(0.5), Inches(1.6), Inches(6), Inches(0.4),
                "EVIDENCE", size=12, bold=True, color=MUTED)
    add_bullets(s, Inches(0.5), Inches(1.95), Inches(6), Inches(3.5), [
        "63-step recorded walkthrough (5 min, narrated)",
        "63 task screenshots stored alongside this deck",
        "Live observations of bottom toolbar, AI panel, FE wizard",
        "Confirmed bug: Logical/Physical AI chip selection does not affect generated output",
    ], size=12)

    add_textbox(s, Inches(7.0), Inches(1.6), Inches(6), Inches(0.4),
                "VALIDATION PLAN", size=12, bold=True, color=MUTED)
    add_bullets(s, Inches(7.0), Inches(1.95), Inches(6), Inches(3.5), [
        "Five 30-min moderated tests with Quest customers",
        "Task: 'Create a retail data model and generate SQL'",
        "Measure: time-on-task, error rate, satisfaction (SUS)",
        "Re-test after each phase and compare deltas",
    ], size=12)

    add_textbox(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.4),
                "OWNERS", size=12, bold=True, color=MUTED)
    add_textbox(s, Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.4),
                "Design lead · Front-end engineering · Product manager · QA",
                size=14, color=INK)

    add_textbox(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
                "Questions or follow-ups: include this deck and the demo video for context.",
                size=11, color=MUTED)


# ---------- Build ----------

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_exec_summary(prs)
    slide_methodology(prs)

    # Hick's Law section
    slide_law_intro(
        prs, "Hick's Law", "Reduce the choices on screen",
        "Decision time grows with the number of options. Group, label, or hide what users don't need right now.",
        "39_canvas_with_generated_entities.png",
    )
    slide_finding(
        prs, NAVY, "HICK'S LAW · FINDING 1",
        "Canvas bottom toolbar: 16 unlabeled icons",
        "The bottom toolbar packs in Select, Add Table, Add View, Add Relationships, Add Annotations, Add Drawing Objects, "
        "View Options, Layout Options, Editors, Mart, Tools, Clipboard, Zoom, Mini-Map and more — all as same-sized "
        "icon buttons with no labels.",
        "First-time users have to hover-then-guess every icon. Discovery of Forward Engineering (under 'Tools') took "
        "two test runs to find. Power users adapt; new users churn.",
        "Group the toolbar into 3 labeled clusters — Build · View · Tools — and add tooltips on every icon. Long-term, "
        "consider a command palette (⌘K) for power users.",
        "39_canvas_with_generated_entities.png", "High",
    )
    slide_finding(
        prs, NAVY, "HICK'S LAW · FINDING 2",
        "AI panel hides the most useful controls",
        "The Model Type and Database selectors live behind a tiny sliders icon at the bottom-left of the textarea, "
        "next to a paperclip with no label. Most users never open them.",
        "Users who want to forward-engineer SQL must pre-select Logical/Physical — but the controls are buried, and "
        "even when discovered the choice doesn't propagate (separate finding).",
        "Promote the controls inline — show 'Type: Logical' and 'Database: Microsoft Fabric' as labeled buttons "
        "directly under the prompt. Keep the popover for advanced settings only.",
        "06_open_the_model_type_selector.png", "High",
    )
    slide_finding(
        prs, NAVY, "HICK'S LAW · FINDING 3",
        "Overview offers four equal-weight Quick Actions",
        "'New Model', 'Open Model', 'Open Mart Model', and 'Reverse Engineering' all sit as same-sized cards. "
        "Below them, the AI Model Generator — likely the primary entry point — competes for the same eye level.",
        "All four cards look like primary actions, so the AI Generator (the differentiator) gets dilluted. Telemetry "
        "would likely show low first-click rates for the AI generator card.",
        "Lift the AI Generator to the hero position; collapse 'Open Mart Model' and 'Reverse Engineering' into a "
        "secondary 'More ways to start' link.",
        "01_sign_in_and_land_on_overview.png", "Medium",
    )

    # Fitts's Law section
    slide_law_intro(
        prs, "Fitts's Law", "Make targets bigger, closer, and predictable",
        "Time to acquire a target is a function of its size and distance. Small icons, far-corner controls, and "
        "double-clicks all add cost.",
        "08_choose_logical_physical.png",
    )
    slide_finding(
        prs, ACCENT, "FITTS'S LAW · FINDING 1",
        "Type/Database icons are 28×28 px — too small to find",
        "The sliders icon that opens the Model Type popover is 28×28 px in a 1440-wide viewport. We had to write "
        "automation to click it precisely.",
        "Touch targets should be ≥40 px (Apple HIG, Material). Anything smaller forces precise pointer work and is "
        "barely usable on touch laptops.",
        "Replace icon-only with labeled buttons: '[ Type: Logical ]' and '[ Database: Microsoft Fabric ]'. Larger "
        "click targets, instant comprehension.",
        "06_open_the_model_type_selector.png", "Medium",
    )
    slide_finding(
        prs, ACCENT, "FITTS'S LAW · FINDING 2",
        "Opening a model requires double-click — not single-click",
        "On the Overview Recent Models list, single-click does nothing. The user has to double-click the file "
        "name to open it. Web users expect single-click on a link.",
        "This violates the strongest convention in web UI. Users will read the row, click once, and assume the "
        "app is broken.",
        "Make the file name a single-click link. Reserve double-click for an advanced 'open with default settings' "
        "shortcut, if needed.",
        "01_sign_in_and_land_on_overview.png", "High",
    )
    slide_finding(
        prs, ACCENT, "FITTS'S LAW · FINDING 3",
        "FE wizard close affordance is in the far corner",
        "The 'X' at (1399, 32) — top-right far corner — is the only way to close the Forward Engineering wizard. "
        "It's also small.",
        "Edge-pinned controls force long mouse travel and are easily mis-clicked, especially on ultrawide "
        "displays. Many designers hit it accidentally while reaching for window controls.",
        "Add a labeled 'Close' button next to 'Back' / 'Generate' at the bottom-right. Keep the corner X as a "
        "secondary affordance.",
        "60_sql_preview_rendered.png", "Low",
    )

    # Jakob's Law section
    slide_law_intro(
        prs, "Jakob's Law", "Keep things familiar",
        "Users spend most of their time on other apps. Quest will feel intuitive when its patterns match what "
        "users already expect.",
        "14_submit_the_ai_prompt.png" if (SCREENS / "14_submit_the_ai_prompt.png").exists() else "23_wait_for_entities_grid.png",
    )
    slide_finding(
        prs, GOLD, "JAKOB'S LAW · WIN",
        "AI Model Generator mirrors ChatGPT — strong familiarity",
        "The textarea + paperclip + 'Create' submit button replicates a pattern hundreds of millions of users "
        "already know. Onboarding effort is near zero.",
        "Lean into this. Anywhere you introduce AI, mirror chat conventions: prompt area, attachments, send "
        "button, conversation history.",
        "Keep this pattern; extend it (e.g., for 'Describe changes' inside the modeler conversation panel — "
        "already done, well).",
        "14_submit_the_ai_prompt.png", "Low",
    )
    slide_finding(
        prs, GOLD, "JAKOB'S LAW · FINDING 1",
        "Disabled controls look like enabled ones",
        "On a Logical-only model, the 'Logical' type dropdown next to Download is greyed but still rendered "
        "as a dropdown. Users click it expecting options.",
        "A disabled control that looks active is a lying interface. Users either give up or file bugs.",
        "Replace with read-only text ('Type: Logical') when not editable. Show a help tooltip explaining when "
        "the type CAN be changed.",
        "49_tools_menu_open.png", "Medium",
    )
    slide_finding(
        prs, GOLD, "JAKOB'S LAW · FINDING 2",
        "AI panel chip selection has no effect (silent failure)",
        "Selecting Logical/Physical + Microsoft Fabric on the AI panel before generating shows the chips, but "
        "the AI produces a Logical-only model anyway. Tools menu then hides 'Forward Engineering' without "
        "explanation.",
        "Users assume their input was respected. When the next screen quietly drops their choice, trust in the "
        "AI evaporates. We hit this in our own walkthrough.",
        "Either honor the chip (preferred) or warn explicitly: 'AI Generator currently produces Logical models; "
        "convert in Properties after creation.' Never silently override.",
        "08_choose_logical_physical.png", "High",
    )

    slide_priority_matrix(prs)
    slide_roadmap(prs)
    slide_appendix(prs)

    prs.save(str(OUT))
    print(f"wrote {OUT}")


if __name__ == "__main__":
    main()
